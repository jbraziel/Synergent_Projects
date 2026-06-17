from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE_PATH = "ACH_Auto_Proposal_Template.pptx"
OUTPUT_PATH = "Generated_ACH_Auto_Proposal.pptx"


proposal_data = {
    "{{proposal_name}}": "ACH Auto Loan Recapture Campaign Proposal",
    "{{proposal_date}}": "April 25, 2026",
    "{{creditunion_name}}": "Sample Credit Union",

    "{{target_conversion_rate}}": "2%",
    "{{total_targets}}": "2,047",
    "{{total_targets_line}}": "Total Targets (de-duped by SSN): 2,047",

    "{{conversions}}": "41",
    "{{avg_auto_balance}}": "$18,500",
    "{{amount_refinanced}}": "$758,500",
    "{{first_year_interest}}": "$45,510",
    "{{auto_interest_rate}}": "6.00%",
    "{{auto_term_years}}": "5",
    "{{target_ROI}}": "325%",

    "{{total_targets_2}}": "4,094",
    "{{total_targets_4}}": "8,188",

    "{{campaign_1_cost}}": "$4,950",
    "{{campaign_2_cost}}": "$8,900",
    "{{campaign_2_per_cost}}": "$4,450",
    "{{campaign_4_cost}}": "$16,020",
    "{{campaign_4_per_cost}}": "$4,005",
}


target_segments = [
    (910, "members making an ACH auto loan payment to another financial institution"),
    (472, "members making a recurring ACH payment between $400–$800 to another financial institution"),
    (628, "members who paid off their auto loan in the last 12 months"),
    (245, "members due to pay off their auto loan in the next 12 months"),
    (2047, "checking members who have a loan but no auto loan with the credit union"),
]

credit_card_target_segments = []


campaign_components = [
    "Creative concept, strategy and design",
    "Preliminary data analysis",
    "Custom programmed data extract for mailing",
    "Proofing and testing",
    "Tracking, monitoring and reporting",
    "Mailing preparation and presorting",
    "Content development",
    "Responsive email template development",
    "Digital graphic assets / social media graphics",
    "5.5” x 8.5” full color variable postcards",
    "Unique URL and QR Code redirect for 12 months",
    "Optional call file for personal outreach / follow-up",
]


def build_dynamic_placeholder_values():
    proposal_data["{{campaign_components}}"] = "\n \n".join(
        f"• {item}" for item in campaign_components
    )

def copy_font_style(source_run, target_run):
    target_run.font.name = source_run.font.name
    target_run.font.size = source_run.font.size
    target_run.font.italic = source_run.font.italic
    target_run.font.underline = source_run.font.underline

    if source_run.font.color.type is not None:
        target_run.font.color.rgb = source_run.font.color.rgb


from pptx.util import Pt

def insert_target_segments(text_frame):
    text_frame.clear()

    TARGET_FONT_SIZE = Pt(10)

    for index, (qty, sentence) in enumerate(target_segments):
        p = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()

        p.space_after = Pt(10)

        run1 = p.add_run()
        run1.text = "•  "
        run1.font.name = "Montserrat"
        run1.font.size = TARGET_FONT_SIZE
        run1.font.bold = False

        run2 = p.add_run()
        run2.text = f"{qty:,}"
        run2.font.name = "Montserrat"
        run2.font.size = TARGET_FONT_SIZE
        run2.font.bold = True
        run2.font.color.rgb = RGBColor(60, 131, 148)

        run3 = p.add_run()
        run3.text = f" {sentence}"
        run3.font.name = "Montserrat"
        run3.font.size = TARGET_FONT_SIZE
        run3.font.bold = False

    # Add Total Targets line after bullets
    p = text_frame.add_paragraph()
    p.space_before = Pt(8)
    
    run1 = p.add_run()
    run1.text = "Total Targets (de-duped by SSN): "
    run1.font.name = "Montserrat"
    run1.font.size = Pt(10)
    run1.font.bold = False
    
    run2 = p.add_run()
    run2.text = proposal_data["{{total_targets}}"]
    run2.font.name = "Montserrat"
    run2.font.size = Pt(10)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor(60, 131, 148)

def insert_credit_card_target_segments(text_frame):
    text_frame.clear()

    TARGET_FONT_SIZE = Pt(10)

    for index, (qty, sentence) in enumerate(credit_card_target_segments):
        p = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()

        p.space_after = Pt(10)

        run1 = p.add_run()
        run1.text = "•  "
        run1.font.name = "Montserrat"
        run1.font.size = TARGET_FONT_SIZE

        run2 = p.add_run()
        run2.text = f"{qty:,}"
        run2.font.name = "Montserrat"
        run2.font.size = TARGET_FONT_SIZE
        run2.font.bold = True
        run2.font.color.rgb = RGBColor(60, 131, 148)

        run3 = p.add_run()
        run3.text = f" {sentence}"
        run3.font.name = "Montserrat"
        run3.font.size = TARGET_FONT_SIZE

    p = text_frame.add_paragraph()
    p.space_before = Pt(8)

    run1 = p.add_run()
    run1.text = "Total Targets (de-duped by SSN): "
    run1.font.name = "Montserrat"
    run1.font.size = Pt(10)

    run2 = p.add_run()
    run2.text = proposal_data["{{total_targets}}"]
    run2.font.name = "Montserrat"
    run2.font.size = Pt(10)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor(60, 131, 148)

def replace_placeholders_in_text_frame(text_frame, replacements):
    for paragraph in text_frame.paragraphs:
        full_text = "".join(run.text for run in paragraph.runs)

        if "{{" not in full_text:
            continue

        new_text = full_text

        for placeholder, value in replacements.items():
            new_text = new_text.replace(placeholder, str(value))

        if new_text != full_text:
            if paragraph.runs:
                first_run = paragraph.runs[0]

                for run in paragraph.runs:
                    run.text = ""

                first_run.text = new_text

def text_frame_contains_placeholder(text_frame, placeholder):
    return placeholder in text_frame.text

def replace_placeholders_in_shapes(shapes, replacements):
    for shape in shapes:

        # Normal text boxes
        if shape.has_text_frame:

            if text_frame_contains_placeholder(shape.text_frame, "{{target_segment_summary}}"):
                insert_target_segments(shape.text_frame)

            if text_frame_contains_placeholder(shape.text_frame, "{{credit_card_target_segment_summary}}"):
                insert_credit_card_target_segments(shape.text_frame)

            else:
                replace_placeholders_in_text_frame(
                    shape.text_frame,
                    replacements
                )

        # Tables, just in case
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    replace_placeholders_in_text_frame(cell.text_frame, replacements)

        # GROUPED SHAPES — this is the missing piece
        if hasattr(shape, "shapes"):
            replace_placeholders_in_shapes(shape.shapes, replacements)


def replace_placeholders_in_ppt(prs, replacements):
    # Normal slide content
    for slide in prs.slides:
        replace_placeholders_in_shapes(slide.shapes, replacements)

    # Layout/footer/header content
    for layout in prs.slide_layouts:
        replace_placeholders_in_shapes(layout.shapes, replacements)

    # Master slide content
    for master in prs.slide_masters:
        replace_placeholders_in_shapes(master.shapes, replacements)


def find_unreplaced_placeholders_in_shapes(shapes, remaining):
    for shape in shapes:
        if shape.has_text_frame:
            for item in shape.text.split():
                if "{{" in item and "}}" in item:
                    remaining.add(item)

        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for item in cell.text.split():
                        if "{{" in item and "}}" in item:
                            remaining.add(item)

        if hasattr(shape, "shapes"):
            find_unreplaced_placeholders_in_shapes(shape.shapes, remaining)


def find_unreplaced_placeholders(prs):
    remaining = set()

    for slide in prs.slides:
        find_unreplaced_placeholders_in_shapes(slide.shapes, remaining)

    for layout in prs.slide_layouts:
        find_unreplaced_placeholders_in_shapes(layout.shapes, remaining)

    for master in prs.slide_masters:
        find_unreplaced_placeholders_in_shapes(master.shapes, remaining)

    return sorted(remaining)

def replace_conversion_estimate_sentence(prs):
    final_sentence = (
        "Conversion estimates are based on the credit union’s average auto loan interest rate of "
        f"{proposal_data['{{auto_interest_rate}}']}, average balance of "
        f"{proposal_data['{{avg_auto_balance}}']} and average term of "
        f"{proposal_data['{{auto_term_years}}']} years."
    )

    sentence_start = (
        "Conversion estimates are based on the credit union’s average auto loan interest rate of"
    )

    for slide in prs.slides:
        replace_sentence_in_shapes(slide.shapes, sentence_start, final_sentence)


def replace_sentence_in_shapes(shapes, sentence_start, final_sentence):
    for shape in shapes:

        if shape.has_text_frame:

            for paragraph in shape.text_frame.paragraphs:

                full_text = "".join(run.text for run in paragraph.runs)

                if sentence_start in full_text:

                    # Clear old text
                    for run in paragraph.runs:
                        run.text = ""

                    FONT_SIZE = Pt(8)

                    # Intro
                    run1 = paragraph.runs[0]
                    run1.text = (
                        "Conversion estimates are based on the credit union’s average auto loan interest rate of "
                    )
                    run1.font.name = "Montserrat"
                    run1.font.size = FONT_SIZE
                    run1.font.bold = False

                    # Interest Rate (bold)
                    run2 = paragraph.add_run()
                    run2.text = proposal_data["{{auto_interest_rate}}"]
                    run2.font.name = "Montserrat"
                    run2.font.size = FONT_SIZE
                    run2.font.bold = True

                    # Middle text
                    run3 = paragraph.add_run()
                    run3.text = ", average balance of "
                    run3.font.name = "Montserrat"
                    run3.font.size = FONT_SIZE
                    run3.font.bold = False

                    # Balance (bold)
                    run4 = paragraph.add_run()
                    run4.text = proposal_data["{{avg_auto_balance}}"]
                    run4.font.name = "Montserrat"
                    run4.font.size = FONT_SIZE
                    run4.font.bold = True

                    # Middle text
                    run5 = paragraph.add_run()
                    run5.text = " and average term of "
                    run5.font.name = "Montserrat"
                    run5.font.size = FONT_SIZE
                    run5.font.bold = False

                    # Term (bold)
                    run6 = paragraph.add_run()
                    run6.text = proposal_data["{{auto_term_years}}"]
                    run6.font.name = "Montserrat"
                    run6.font.size = FONT_SIZE
                    run6.font.bold = True

                    # Ending
                    run7 = paragraph.add_run()
                    run7.text = " years."
                    run7.font.name = "Montserrat"
                    run7.font.size = FONT_SIZE
                    run7.font.bold = False

        if hasattr(shape, "shapes"):
            replace_sentence_in_shapes(
                shape.shapes,
                sentence_start,
                final_sentence
            )

def replace_total_targets_line(prs):
    final_line = (
        "Total Targets (de-duped by SSN): "
        f"{proposal_data['{{total_targets}}']}"
    )

    for slide in prs.slides:
        replace_total_targets_in_shapes(slide.shapes, final_line)


def replace_total_targets_in_shapes(shapes, final_line):
    for shape in shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs)

                if "{{total_targets_line}}" in full_text:
                    if paragraph.runs:
                        first_run = paragraph.runs[0]

                        for run in paragraph.runs:
                            run.text = ""

                        first_run.text = final_line
                        first_run.font.name = "Montserrat"
                        first_run.font.size = Pt(8)
                        first_run.font.bold = False

        if hasattr(shape, "shapes"):
            replace_total_targets_in_shapes(shape.shapes, final_line)

def format_conversion_objective_line(prs):
    sentence_start = "Deliver Measurable Campaign Performance"

    for slide in prs.slides:
        format_conversion_objective_in_shapes(slide.shapes, sentence_start)


def format_conversion_objective_in_shapes(shapes, sentence_start):
    for shape in shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs)

                if sentence_start in full_text:
                    for run in paragraph.runs:
                        run.text = ""

                    run1 = paragraph.runs[0]
                    run1.text = "Deliver Measurable Campaign Performance "
                    run1.font.name = "Montserrat"
                    run1.font.size = Pt(10.5)
                    run1.font.bold = True

                    run2 = paragraph.add_run()
                    run2.text = "by attaining a "
                    run2.font.name = "Montserrat"
                    run2.font.size = Pt(10.5)
                    run2.font.bold = False

                    run3 = paragraph.add_run()
                    run3.text = proposal_data["{{target_conversion_rate}}"]
                    run3.font.name = "Montserrat"
                    run3.font.size = Pt(10.5)
                    run3.font.bold = True
                    run3.font.color.rgb = RGBColor(60, 131, 148)

                    run4 = paragraph.add_run()
                    run4.text = (
                        " conversion rate within the targeted segment, translating into increased loan originations, "
                        "recaptured balances, and improved return on marketing investment."
                    )
                    run4.font.name = "Montserrat"
                    run4.font.size = Pt(10.5)
                    run4.font.bold = False

        if hasattr(shape, "shapes"):
            format_conversion_objective_in_shapes(shape.shapes, sentence_start)

def format_emp_page_10_costs(prs):
    BLUE_GREEN = RGBColor(60, 131, 148)
    FONT_SIZE = Pt(8)

    emp_values = [
        proposal_data["{{tier_cost}}"],
        proposal_data["{{essentials_cost}}"],
        proposal_data["{{premium_cost}}"],
        proposal_data["{{elite_cost}}"],
    ]

    if len(prs.slides) < 10:
        return

    slide = prs.slides[9]
    format_emp_page_10_costs_in_shapes(slide.shapes, emp_values, BLUE_GREEN, FONT_SIZE)


def format_emp_page_10_costs_in_shapes(shapes, emp_values, BLUE_GREEN, FONT_SIZE):
    for shape in shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs).strip()

                # Only package boxes, not Custom Set-up
                if not full_text.startswith("Monthly Cost:"):
                    continue

                for value in emp_values:
                    if value in full_text:
                        for run in paragraph.runs:
                            run.text = ""

                        first_run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        first_run.text = "Monthly Cost: "
                        first_run.font.name = "Montserrat"
                        first_run.font.size = FONT_SIZE
                        first_run.font.bold = False

                        value_run = paragraph.add_run()
                        value_run.text = value
                        value_run.font.name = "Montserrat"
                        value_run.font.size = FONT_SIZE
                        value_run.font.bold = True
                        value_run.font.color.rgb = BLUE_GREEN

                        break

        if hasattr(shape, "shapes"):
            format_emp_page_10_costs_in_shapes(shape.shapes, emp_values, BLUE_GREEN, FONT_SIZE)

def highlight_emp_pricing_tier(prs):
    GREEN = RGBColor(118, 189, 34)
    WHITE = RGBColor(255, 255, 255)

    tier_number = str(proposal_data.get("{{emp_tier_number}}", "")).strip()
    print("DEBUG EMP tier number:", tier_number)

    if not tier_number:
        return

    # Page 8 = slide index 7
    if len(prs.slides) < 8:
        return

    slide = prs.slides[7]

    for shape in slide.shapes:
        if not shape.has_table:
            continue

        table = shape.table

        for row in table.rows:
            row_text = " ".join(cell.text.strip() for cell in row.cells)

            # Match the tier row
            if row_text.startswith(tier_number + " "):
                for cell in row.cells:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = GREEN

                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = WHITE
                            run.font.bold = True
                            run.font.size = Pt(10.5)


def main(output_path=None):
    build_dynamic_placeholder_values()

    prs = Presentation(TEMPLATE_PATH)

    replace_placeholders_in_ppt(prs, proposal_data)
    format_conversion_objective_line(prs)
    replace_total_targets_line(prs)
    replace_conversion_estimate_sentence(prs)

    is_emp_proposal = "{{tier_cost}}" in proposal_data and "{{essentials_cost}}" in proposal_data
    if is_emp_proposal:
      format_emp_page_10_costs(prs)
      highlight_emp_pricing_tier(prs)
    
    if output_path is None:
        output_path = OUTPUT_PATH

    prs.save(output_path)

    print(f"Proposal created successfully: {output_path}")

    remaining = find_unreplaced_placeholders(prs)

    if remaining:
        print("\nThese placeholders may still be unreplaced:")
        for item in remaining:
            print(f" - {item}")
    else:
        print("All placeholders were replaced.")


if __name__ == "__main__":
    main()
